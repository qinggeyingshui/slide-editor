"""
导出 slide-editor 为 PPTX（每页截图为全屏图片插入）
用法: python tools/export_pptx.py [--url URL] [--output PATH]
默认: url=http://localhost:3001  output=./exported.pptx

原理: Selenium打开编辑器→逐页截图canvas区域→python-pptx组装
"""
import argparse, time, os, tempfile, shutil, base64
from pathlib import Path

def export_slides_to_pptx(url="http://localhost:3001", output="exported.pptx"):
    from selenium import webdriver
    from selenium.webdriver.chrome.options import Options
    from pptx import Presentation
    from pptx.util import Inches, Emu

    # Launch headless Chrome
    opts = Options()
    opts.add_argument("--headless=new")
    opts.add_argument("--window-size=1920,1080")
    opts.add_argument("--force-device-scale-factor=2")
    opts.add_argument("--disable-gpu")
    opts.add_argument("--no-sandbox")
    driver = webdriver.Chrome(options=opts)

    tmp_dir = tempfile.mkdtemp(prefix="slide_export_")
    screenshots = []

    try:
        driver.get(url)
        time.sleep(3)

        # Get slide count from thumbnails
        slide_count = driver.execute_script(
            "return document.querySelectorAll('.slide-thumb').length"
        )
        if not slide_count:
            slide_count = 9
        print(f"Found {slide_count} slides")

        # Screenshot each slide
        for i in range(slide_count):
            # Click thumbnail to switch slide
            driver.execute_script(f"""
                const thumbs = document.querySelectorAll('.slide-thumb');
                if (thumbs[{i}]) thumbs[{i}].click();
            """)
            time.sleep(0.5)

            # Screenshot the canvas area
            canvas = driver.execute_script("""
                const el = document.querySelector('.slide-render');
                if (!el) return null;
                const rect = el.getBoundingClientRect();
                return {x: rect.x, y: rect.y, w: rect.width, h: rect.height};
            """)
            if not canvas:
                print(f"  Slide {i+1}: canvas not found, using full page")
                path = os.path.join(tmp_dir, f"slide_{i:02d}.png")
                driver.save_screenshot(path)
            else:
                # Use CDP to screenshot specific region
                result = driver.execute_cdp_cmd("Page.captureScreenshot", {
                    "format": "png",
                    "clip": {
                        "x": canvas['x'], "y": canvas['y'],
                        "width": canvas['w'], "height": canvas['h'],
                        "scale": 2
                    }
                })
                path = os.path.join(tmp_dir, f"slide_{i:02d}.png")
                with open(path, 'wb') as f:
                    f.write(base64.b64decode(result['data']))

            screenshots.append(path)
            print(f"  Slide {i+1}/{slide_count} captured")


        # Assemble PPTX
        prs = Presentation()
        # Set 16:9 slide size
        prs.slide_width = Emu(12192000)   # 10 inches * 914400
        prs.slide_height = Emu(6858000)   # 7.5 inches * 914400

        blank_layout = prs.slide_layouts[6]  # blank layout

        for img_path in screenshots:
            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(
                img_path, Emu(0), Emu(0),
                width=prs.slide_width, height=prs.slide_height
            )

        out_path = os.path.abspath(output)
        os.makedirs(os.path.dirname(out_path) if os.path.dirname(out_path) else '.', exist_ok=True)
        prs.save(out_path)
        print(f"\nDone! Exported {slide_count} slides to: {out_path}")
        return out_path

    finally:
        driver.quit()
        shutil.rmtree(tmp_dir, ignore_errors=True)


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Export slide-editor to PPTX")
    parser.add_argument("--url", default="http://localhost:3001")
    parser.add_argument("--output", "-o", default="exported.pptx")
    args = parser.parse_args()
    export_slides_to_pptx(args.url, args.output)
